# -*- coding: utf-8 -*-
import io, os
from pptx import Presentation
PROJ = r"C:\Users\jee\University\3rd grade\Data_visualization\lol_dataviz"
f = os.path.join(PROJ, "발표자료_롤프로지표.pptx")
report = os.path.join(PROJ, "qa_report.txt")
prs = Presentation(f)
SW, SH = prs.slide_width, prs.slide_height
def inch(e): return round(e/914400, 2)
out = io.StringIO()
oob_count = 0
out.write(f"slide size: {inch(SW)} x {inch(SH)}\n")
for i, s in enumerate(prs.slides, 1):
    pics = sum(1 for sh in s.shapes if sh.shape_type == 13)
    out.write(f"\n=== Slide {i}  (pictures={pics}) ===\n")
    for sh in s.shapes:
        oob = ""
        l,t,w,h = sh.left, sh.top, sh.width, sh.height
        if None not in (l,t,w,h):
            if l < -9144 or t < -9144 or (l+w) > SW+9144 or (t+h) > SH+9144:
                oob = f"  <<< OOB l={inch(l)} t={inch(t)} r={inch(l+w)} b={inch(t+h)}"
                oob_count += 1
        txt = ""
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt = " | " + sh.text_frame.text.strip().replace("\n", " / ")
        out.write(f"  type={sh.shape_type}{txt}{oob}\n")
with open(report, "w", encoding="utf-8") as fh:
    fh.write(out.getvalue())
print("OOB_COUNT", oob_count)
print("report written:", report)
