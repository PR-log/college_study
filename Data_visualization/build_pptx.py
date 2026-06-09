# -*- coding: utf-8 -*-
"""데이터로 보는 프로 LoL — 발표자료(.pptx) 생성 (python-pptx)"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

PROJ = r"C:\Users\jee\University\3rd grade\Data_visualization\lol_dataviz"
OUT_DIR = os.path.join(PROJ, "output")
OUTFILE = os.path.join(PROJ, "발표자료_롤프로지표.pptx")
FONT = "Malgun Gothic"

# ---- palette (LoL 톤) ----
NAVY = RGBColor(0x0A, 0x14, 0x28)
GOLD = RGBColor(0xC8, 0xAA, 0x6E)
TEAL = RGBColor(0x0A, 0xC8, 0xB9)
INK  = RGBColor(0x1E, 0x23, 0x28)
MUTE = RGBColor(0x6B, 0x76, 0x80)
CARD = RGBColor(0xF1, 0xF4, 0xF6)
WHITE= RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def _setfont(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)

def text(s, l, t, w, h, paras, anchor=MSO_ANCHOR.TOP):
    """paras: list of dict(text,size,color,bold,italic,align,space_after)"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, tf.margin_right): pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("space_after") is not None:
            para.space_after = Pt(p["space_after"])
        if p.get("space_before") is not None:
            para.space_before = Pt(p["space_before"])
        r = para.add_run(); r.text = p["text"]
        r.font.size = Pt(p["size"]); r.font.bold = p.get("bold", False)
        r.font.italic = p.get("italic", False)
        r.font.color.rgb = p.get("color", INK)
        _setfont(r)
    return tb

def rect(s, l, t, w, h, fill, line=None, round_=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def fit(path, bl, bt, bw, bh):
    iw, ih = Image.open(path).size
    r = iw / ih
    w = bw; h = w / r
    if h > bh:
        h = bh; w = h * r
    return bl + (bw - w) / 2, bt + (bh - h) / 2, w, h

def image(s, name, bl, bt, bw, bh):
    p = os.path.join(OUT_DIR, name)
    l, t, w, h = fit(p, bl, bt, bw, bh)
    s.shapes.add_picture(p, Inches(l), Inches(t), Inches(w), Inches(h))

def title(s, txt):
    text(s, 0.7, 0.45, 12.0, 0.85,
         [dict(text=txt, size=28, bold=True, color=NAVY)])

def pagenum(s, n):
    text(s, 12.4, 7.0, 0.7, 0.35,
         [dict(text=str(n), size=11, color=MUTE, align=PP_ALIGN.RIGHT)])

def insight(s, l, t, w, h, head, body):
    rect(s, l, t, w, h, NAVY, round_=True)
    tb = s.shapes.add_textbox(Inches(l+0.3), Inches(t+0.25), Inches(w-0.6), Inches(h-0.5))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p0 = tf.paragraphs[0]; p0.space_after = Pt(8)
    r = p0.add_run(); r.text = head; r.font.size=Pt(17); r.font.bold=True; r.font.color.rgb=GOLD; _setfont(r)
    p1 = tf.add_paragraph(); r = p1.add_run(); r.text = body
    r.font.size=Pt(14); r.font.color.rgb=WHITE; _setfont(r)

# ===================== Slide 1 : 표지 =====================
s = slide(NAVY)
# 모티프: 골드/틸 작은 사각형
for i, c in enumerate((GOLD, TEAL, GOLD)):
    rect(s, 0.9, 2.35 + i*0.34, 0.2, 0.2, c)
text(s, 1.35, 2.1, 11.2, 3.0, [
    dict(text="데이터로 보는 프로 LoL", size=46, bold=True, color=GOLD, space_after=6),
    dict(text="무엇이 강한 선수·강한 팀을 만드는가", size=24, color=WHITE, space_after=18),
    dict(text="데이터시각화  ·  Oracle's Elixir (2024–2025)  ·  R / ggplot2", size=14, color=TEAL),
])
text(s, 1.35, 6.35, 11.0, 0.5,
     [dict(text="이름 ○○○   |   학번 000000   |   데이터시각화", size=13, color=MUTE)])

# ===================== Slide 2 : 동기 =====================
s = slide()
title(s, "왜 이 주제인가")
text(s, 0.7, 2.0, 6.7, 4.0, [
    dict(text="저는 평소 롤을 즐겨 하고", size=18, color=INK, space_after=4),
    dict(text="프로 경기도 자주 봅니다.", size=18, color=INK, space_after=18),
    dict(text="“그냥 잘한다”는 느낌이 아니라,", size=18, color=INK, space_after=4),
    dict(text="강한 선수의 차이를 실제 데이터로", size=18, color=INK, space_after=4),
    dict(text="확인해보고 싶어 이 주제를 골랐습니다.", size=18, color=INK),
])
# 우측 질문 카드
rect(s, 7.9, 2.0, 4.7, 3.4, NAVY, round_=True)
tb = s.shapes.add_textbox(Inches(8.2), Inches(2.0), Inches(4.1), Inches(3.4))
tf = tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
tf.margin_left=0; tf.margin_right=0
p = tf.paragraphs[0]; r=p.add_run(); r.text="“"; r.font.size=Pt(40); r.font.bold=True; r.font.color.rgb=GOLD; _setfont(r)
p2 = tf.add_paragraph(); r=p2.add_run(); r.text="강한 선수는\n대체 뭐가 다를까?"
r.font.size=Pt(26); r.font.bold=True; r.font.color.rgb=WHITE; _setfont(r)
pagenum(s, 2)

# ===================== Slide 3 : 데이터 & 방법 =====================
s = slide()
title(s, "데이터 & 방법")
cards = [("데이터 출처","Oracle's Elixir"),("분석 기간","2024 – 2025 시즌"),
         ("데이터 규모","선수 기록 25,350건"),("분석 도구","R · ggplot2")]
cw, ch = 5.75, 1.95
pos = [(0.7,1.7),(6.85,1.7),(0.7,3.85),(6.85,3.85)]
for (lab,val),(cl,ct) in zip(cards,pos):
    rect(s, cl, ct, cw, ch, CARD, round_=True)
    tb = s.shapes.add_textbox(Inches(cl+0.35), Inches(ct+0.3), Inches(cw-0.7), Inches(ch-0.6))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=0; tf.margin_right=0
    p=tf.paragraphs[0]; p.space_after=Pt(6); r=p.add_run(); r.text=lab
    r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=RGBColor(0xA8,0x8A,0x4E); _setfont(r)
    p2=tf.add_paragraph(); r=p2.add_run(); r.text=val
    r.font.size=Pt(23); r.font.bold=True; r.font.color.rgb=NAVY; _setfont(r)
text(s, 0.7, 6.15, 12.0, 0.7,
     [dict(text="※ 한계: 중국 LPL 정규시즌은 상세 지표가 공개되지 않아, 국제대회(MSI·Worlds) 경기만 포함됩니다.",
           size=13, color=MUTE)])
pagenum(s, 3)

# ===================== Slide 4 : 포지션별 딜량 =====================
s = slide()
title(s, "① 포지션마다 중요한 지표가 다르다")
image(s, "01_position_dpm.png", 0.5, 1.5, 7.9, 5.3)
insight(s, 8.7, 2.6, 4.0, 2.9,
        "역할마다 다른 기준",
        "바텀·미드는 딜량이 핵심, 서포터는 가장 낮음. 같은 ‘잘함’도 포지션마다 평가 기준이 달라야 한다.")
pagenum(s, 4)

# ===================== Slide 5 : 15분 골드 격차 =====================
s = slide()
title(s, "② 초반 우위가 승패를 가른다")
image(s, "02_golddiff15_winloss.png", 0.5, 1.5, 7.9, 5.3)
insight(s, 8.7, 2.6, 4.0, 2.9,
        "15분이 결정적",
        "이긴 경기(파랑)는 모든 라인에서 15분 골드가 앞선다. 후반 한방보다 초반 우위가 승패와 직결.")
pagenum(s, 5)

# ===================== Slide 6 : 딜량 vs 승률 =====================
s = slide()
title(s, "③ 핵심 지표는 실제 승률로 이어진다")
image(s, "03_player_dpm_winrate.png", 0.6, 1.55, 12.1, 4.55)
text(s, 0.7, 6.35, 12.0, 0.7,
     [dict(text="모든 포지션에서 우상향 — 평균 딜량이 높은 선수일수록 승률도 높다.",
           size=16, bold=True, color=NAVY)])
pagenum(s, 6)

# ===================== Slide 7 : 최고의 선수 =====================
s = slide()
title(s, "④ 그래서, 최고의 미드는?")
image(s, "04_top10_mid_dpm.png", 0.4, 1.5, 6.6, 4.6)
image(s, "05_radar_mid.png", 7.1, 1.4, 5.7, 4.8)
text(s, 0.6, 6.45, 12.1, 0.7,
     [dict(text="딜량 1위 Chovy — 딜·CS·골드·시야·KDA 5개 지표를 모두 압도하는 ‘약점 없는 선수’.",
           size=16, bold=True, color=NAVY)])
pagenum(s, 7)

# ===================== Slide 8 : 메타 변화 =====================
s = slide()
title(s, "⑤ 게임은 더 빨라졌다")
image(s, "06_meta_gamelength.png", 0.5, 1.5, 7.9, 5.3)
insight(s, 8.7, 2.6, 4.0, 2.9,
        "32분대로 단축·안정화",
        "평균 경기 시간이 2016년 38분 정점 이후 줄어 2019년부터 32분 안팎. 패치가 ‘빠른 게임’을 지향.")
pagenum(s, 8)

# ===================== Slide 9 : 결론 =====================
s = slide(NAVY)
text(s, 0.9, 0.7, 11.5, 0.9, [dict(text="결론", size=34, bold=True, color=GOLD)])
text(s, 0.9, 1.9, 11.5, 1.4, [
    dict(text="강한 선수 = 포지션 핵심 지표 상위  +  초반 우위 창출", size=24, bold=True, color=WHITE),
])
benefits = [("솔랭 유저","포지션별로 무엇을 연습할지 — 연습 우선순위"),
            ("팬","경기를 보는 새로운 데이터 관점"),
            ("코치·분석가","객관적인 선수 평가 기준(KPI)")]
ty = 3.4
for head, body in benefits:
    rect(s, 0.95, ty+0.09, 0.18, 0.18, GOLD)
    tb = s.shapes.add_textbox(Inches(1.35), Inches(ty-0.03), Inches(11.2), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = head + "  →  "
    r1.font.size=Pt(17); r1.font.bold=True; r1.font.color.rgb=GOLD; _setfont(r1)
    r2 = p.add_run(); r2.text = body
    r2.font.size=Pt(16); r2.font.color.rgb=WHITE; _setfont(r2)
    ty += 0.72
text(s, 0.95, 6.5, 11.5, 0.6,
     [dict(text="막연한 ‘잘한다’를 구체적인 지표로.", size=14, italic=True, color=TEAL)])

prs.save(OUTFILE)
print("SAVED:", OUTFILE)
print("slides:", len(prs.slides._sldIdLst))
